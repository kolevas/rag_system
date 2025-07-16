from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.table import Table
from pptx.shapes.autoshape import Shape
from typing import List, Dict, Tuple
import re
import unicodedata
import time
from pathlib import Path
import io

def clean_text(text: str) -> str:
    """
    Clean and normalize text by removing extra whitespace, special characters,
    and normalizing unicode characters. Preserves table formatting.
    """
    # print("Starting text cleaning...")
    
    # Split text into table and non-table parts
    parts = text.split("\nTable Content:")
    cleaned_parts = []
    
    # Clean the first part (non-table content)
    if parts[0].strip():
        # print("Cleaning non-table content...")
        # Normalize unicode characters
        cleaned = unicodedata.normalize('NFKD', parts[0])
        # print("Unicode normalization completed")
        
        # Remove special characters but keep basic punctuation
        cleaned = re.sub(r'[^\w\s.,!?-]', ' ', cleaned)
        # print("Special characters removed")
        
        # Remove extra whitespace
        cleaned = re.sub(r'\s+', ' ', cleaned)
        # print("Extra whitespace removed")
        
        cleaned_parts.append(cleaned.strip())
    
    # Process remaining parts (tables)
    for part in parts[1:]:
        if part.strip():
            # print("Preserving table content...")
            # Add back the "Table Content:" marker
            cleaned_parts.append("\nTable Content:" + part)
    
    return '\n'.join(cleaned_parts)

def process_table(table: Table) -> str:
    """
    Process a PowerPoint table and return its content in a structured text format.
    """
    print(f"Processing table with {len(table.rows)} rows...")
    table_text = []
    for row_idx, row in enumerate(table.rows):
        row_text = []
        for cell in row.cells:
            # Clean cell text and add it to row
            cell_text = clean_text(cell.text)
            if cell_text:
                row_text.append(cell_text)
        if row_text:
            table_text.append(" | ".join(row_text))
        print(f"Processed row {row_idx + 1}/{len(table.rows)}")
    final_output =  "\n".join(table_text)

    return final_output.join("End of table content.\n\n")

def extract_text_from_slide(slide: Slide) -> List[Tuple[str, str]]:
    """
    Extract text from a PowerPoint slide, including tables and shapes.
    Returns a list of tuples (content_type, content) to preserve order.
    """
    content_parts = []
    
    # Process slide title
    if slide.shapes.title:
        title_text = slide.shapes.title.text.strip()
        if title_text:
            content_parts.append(('title', f" {title_text}"))
    
    # Process all shapes in the slide
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            if isinstance(shape, Table):
                table_text = process_table(shape)
                if table_text:
                    content_parts.append(('table', table_text))
            else:
                # Regular text shape
                content_parts.append(('text', shape.text))
    
    return content_parts

def extract_text_from_pptx(file) -> List[Tuple[str, str]]:
    """
    Extract text from a PowerPoint presentation, including all slides and their content.
    Returns a list of tuples (content_type, content) to preserve order.
    """
    try:
        print(f"\nOpening presentation: {file["name"]}")
        prs = Presentation(io.BytesIO(file["content"]))
        all_content = []
        
        print("\nProcessing slides in order...")
        for slide_idx, slide in enumerate(prs.slides):
            # print(f"Processing slide {slide_idx + 1}/{len(prs.slides)}")
            slide_content = extract_text_from_slide(slide)
            if slide_content:
                all_content.extend(slide_content)
        
        print("\nPresentation processing completed")
        return all_content
    except Exception as e:
        print(f"\nError during presentation processing: {str(e)}")
        raise Exception(f"Error extracting text from presentation: {str(e)}")

def split_text_into_chunks(text: str, chunk_size: int = 1500, overlap: int = 200) -> List[str]:
    """
    Split text into overlapping chunks of specified size.
    - If a table starts in a chunk and cannot be stored as a whole, move it to the next chunk.
    - If a table is larger than the chunk size, split by rows, each chunk starts with the header.
    - If a table fits in the chunk, keep it with other data in the same chunk.
    """
    chunks = []
    text_length = len(text)
    start = 0

    min_overlap = max(50, int(overlap * 0.5)) # Ensure a minimum overlap, at least 50
    max_overlap = min(overlap * 1.5, chunk_size - 50) # Max overlap should not be too close to chunk_size

    # Find all tables
    table_matches = list(re.finditer(r"(Table Content:.*?End of table content\.)", text, re.DOTALL))
    table_spans = [m.span() for m in table_matches]

    while start < text_length:
        # Find the next table that starts after 'start'
        next_table = None
        for t_start, t_end in table_spans:
            if t_start >= start:
                next_table = (t_start, t_end)
                break

        end = start + chunk_size
        if end > text_length:
            end = text_length

        if next_table and next_table[0] < end:
            t_start, t_end = next_table
            # If the table fits entirely within the chunk, include it with other data
            if t_end <= end:
                # Find proper sentence boundary before the table or at the end
                best_end = end
                if end < text_length:
                    # Look for sentence boundaries first (period, exclamation, question mark)
                    sentence_end = -1
                    for i in range(end, max(start + chunk_size // 2, start), -1):
                        if i < text_length and text[i-1] in '.!?' and (i == text_length or text[i].isspace() or text[i].isupper()):
                            sentence_end = i
                            break
                    
                    if sentence_end != -1:
                        best_end = sentence_end
                    else:
                        # If no sentence boundary found, fall back to other boundaries
                        for i in range(end, max(start + chunk_size // 2, start), -1):
                            if i < text_length and text[i-1] in '.\n':
                                best_end = i
                                break
                
                chunk = text[start:best_end]
                chunks.append(chunk)
                print(f"Created chunk {len(chunks)} (length: {len(chunk)})")
                if best_end >= text_length:
                    break
                
                # Find where to start the next chunk (after any whitespace)
                next_start = best_end
                while next_start < text_length and text[next_start].isspace():
                    next_start += 1
                
                start = next_start
            else:
                # Table would be split, so move it to the next chunk
                if t_start > start:
                    # Add text before the table as a chunk (use original logic)
                    pre_table_text = text[start:t_start]
                    if pre_table_text.strip():
                        pre_chunks = _original_chunking(pre_table_text, chunk_size, overlap)
                        chunks.extend(pre_chunks)
                # Now handle the table
                if t_end - t_start > chunk_size:
                    # Table is too big, split by rows with header
                    table_text = text[t_start:t_end]
                    table_lines = table_text.split('\n')
                    header = table_lines[1].strip() if len(table_lines) > 1 else ''
                    rows = table_lines[2:-1]
                    current_chunk = f"Table Content:\n{header}"
                    for row in rows:
                        if len(current_chunk) + len(row) + 1 > chunk_size:
                            chunks.append(current_chunk)
                            print(f"Created chunk {len(chunks)} (length: {len(current_chunk)})")
                            current_chunk = f"Table Content:\n{header}\n{row}"
                        else:
                            current_chunk += f"\n{row}"
                    if current_chunk:
                        chunks.append(current_chunk + "\nEnd of table content.")
                        print(f"Created chunk {len(chunks)} (length: {len(current_chunk)})")
                else:
                    # Table fits as a whole, add it as a chunk
                    chunks.append(text[t_start:t_end])
                    print(f"Created chunk {len(chunks)} (length: {len(text[t_start:t_end])})")
                start = t_end
        else:
            # No table in this chunk, find proper sentence boundaries
            end = start + chunk_size
            if end >= text_length:
                end = text_length
                chunk = text[start:end]
                chunks.append(chunk)
                break
            else:
                # Find the best place to end the chunk at a sentence boundary
                best_end = end
                # Look backwards from the target end to find a sentence ending
                sentence_end = -1
                for i in range(end, max(start + chunk_size // 2, start), -1):
                    if i < text_length and text[i-1] in '.!?' and (i == text_length or text[i].isspace() or text[i].isupper()):
                        sentence_end = i
                        break
                
                if sentence_end != -1:
                    best_end = sentence_end
                else:
                    # If no sentence boundary found, look for other boundaries
                    for i in range(end, max(start + chunk_size // 2, start), -1):
                        if i < text_length and text[i-1] in '.\n':
                            best_end = i
                            break
                
                chunk = text[start:best_end]
                chunks.append(chunk)
                print(f"Created chunk {len(chunks)} (length: {len(chunk)})")
                
                # Find where to start the next chunk (after any whitespace)
                next_start = best_end
                while next_start < text_length and text[next_start].isspace():
                    next_start += 1
                
                start = next_start

    print(f"Text splitting completed. Total chunks: {len(chunks)}")
    return chunks

# Helper function for original chunking logic

def _original_chunking(text: str, chunk_size: int, overlap: int) -> List[str]:
    chunks = []
    text_length = len(text)
    start = 0
    min_overlap = max(50, int(overlap * 0.5)) # Ensure a minimum overlap, at least 50
    max_overlap = min(overlap * 1.5, chunk_size - 50) # Max overlap should not be too close to chunk_size
    while start < text_length:
        end = start + chunk_size
        if end < text_length:
            # Look for sentence boundaries first (period, exclamation, question mark)
            sentence_end = -1
            for i in range(end, start, -1):
                if text[i-1] in '.!?' and i < len(text) and (text[i].isspace() or text[i].isupper()):
                    sentence_end = i
                    break
            
            # If no sentence boundary found, fall back to other boundaries
            if sentence_end == -1:
                last_period = text.rfind('.', start, end)
                last_exclamation = text.rfind('!', start, end)
                last_question = text.rfind('?', start, end)
                last_newline = text.rfind('\n', start, end)
                last_space = text.rfind(' ', start, end)
                split_point = max(last_period, last_exclamation, last_question, last_newline, last_space)
            else:
                split_point = sentence_end
            if split_point > start:
                end = split_point + 1
        else:
            end = text_length
        chunk = text[start:end]
        chunks.append(chunk)
        print(f"Created chunk {len(chunks)} (length: {len(chunk)})")
        if end == text_length:
            break
        potential_overlap_start = max(start, end - overlap)
        overlap_candidate = -1
        for i in range(end - min_overlap, potential_overlap_start -1, -1):
            if text[i] in '.!?' and i + 1 < len(text) and (text[i + 1].isspace() or text[i + 1].isupper()):
                overlap_candidate = i + 1  # Start after the sentence end
                break
            elif text[i] == '.' or text[i] == ' ' or text[i] == '\n':
                overlap_candidate = i + 1
                break
        current_overlap = overlap
        if overlap_candidate != -1:
            calculated_overlap = end - overlap_candidate
            if min_overlap <= calculated_overlap <= max_overlap:
                current_overlap = calculated_overlap
            elif calculated_overlap < min_overlap:
                current_overlap = min_overlap
            else:
                current_overlap = max_overlap
        else:
            current_overlap = max(min_overlap, min(overlap, max_overlap))
            print(f"No strategic overlap point found. Using default/clamped overlap: {current_overlap}")
        start = end - current_overlap
    return chunks

def normalize_chunk(chunk: str) -> str:
    """
    Normalize a text chunk by cleaning up formatting, spacing, and punctuation.
    """
    normalized_response = re.sub(r"\[\[\d+\]\],?\s*", "", chunk)
    normalized_response = re.sub(r"\n\s*#{1,6}\s*([^\n]+)", r". \1:", normalized_response)
    normalized_response = re.sub(r"\n\s*-\s*\*\*([^*]+)\*\*:\s*", r". \1: ", normalized_response)
    normalized_response = re.sub(r"\n\s*-\s*", ". ", normalized_response)
    normalized_response = re.sub(r"\*\*([^*]+)\*\*", r"\1", normalized_response)
    normalized_response = re.sub(r"\*([^*]+)\*", r"\1", normalized_response)
    normalized_response = re.sub(r"\n{3,}", "\n\n", normalized_response)
    normalized_response = re.sub(r"\n\n", ". ", normalized_response)
    normalized_response = re.sub(r"\n", " ", normalized_response)
    
    # Handle hyphenated words broken across lines
    normalized_response = re.sub(r"(\w)-\s+(\w)", r"\1\2", normalized_response)
    
    # Remove extra hyphens and normalize dashes
    normalized_response = re.sub(r"--+", "—", normalized_response)
    
    # Fix common PDF extraction artifacts
    normalized_response = re.sub(r"\s*\|\s*", " ", normalized_response)  # Remove table separators
    normalized_response = re.sub(r"(\w)\s*_\s*(\w)", r"\1_\2", normalized_response)  # Fix broken underscores
    
    # Handle page numbers and headers/footers patterns
    normalized_response = re.sub(r"\b(Page\s+\d+|©\s*\d+|All rights reserved)\b", "", normalized_response, flags=re.IGNORECASE)
    
    # Fix spacing around parentheses and brackets
    normalized_response = re.sub(r"\s*\(\s*", " (", normalized_response)
    normalized_response = re.sub(r"\s*\)\s*", ") ", normalized_response)
    normalized_response = re.sub(r"\s*\[\s*", " [", normalized_response)
    normalized_response = re.sub(r"\s*\]\s*", "] ", normalized_response)
    
    # Clean up punctuation
    normalized_response = re.sub(r"\.+", ".", normalized_response)
    normalized_response = re.sub(r"\s*\.\s*\.", ".", normalized_response)
    normalized_response = re.sub(r":\s*\.", ":", normalized_response)
    normalized_response = re.sub(r"\.\s*:", ":", normalized_response)
    
    # Clean up multiple punctuation marks
    normalized_response = re.sub(r"[.]{2,}", ".", normalized_response)
    normalized_response = re.sub(r"[,]{2,}", ",", normalized_response)
    
    normalized_response = re.sub(r"\s+", " ", normalized_response)
    normalized_response = re.sub(r"\s*([.,:;!?])", r"\1", normalized_response)
    normalized_response = re.sub(r"([.,:;!?])\s*", r"\1 ", normalized_response)
    
    # Fix spacing before opening quotes and after closing quotes
    normalized_response = re.sub(r'\s*"([^"]*)"\s*', r' "\1" ', normalized_response)
    normalized_response = re.sub(r"\s*'([^']*)'\s*", r" '\1' ", normalized_response)
    
    normalized_response = re.sub(r"\.\s*([a-z])", lambda m: ". " + m.group(1).upper(), normalized_response)
    normalized_response = re.sub(r"(\d+)\.\s+(\d+)", r"\1.\2", normalized_response)
    
    # Final cleanup - remove any remaining double spaces
    normalized_response = re.sub(r"\s{2,}", " ", normalized_response)
    
    normalized_response = normalized_response.strip()
    if normalized_response and not normalized_response[0].isupper():
        normalized_response = normalized_response[0].upper() + normalized_response[1:]
    if normalized_response and normalized_response[-1] not in '.!?':
        normalized_response += "."
    return normalized_response

def preprocess_presentation(file, chunk_size: int = 1000, overlap: int = 200, blob_metadata=None) -> Dict[str, List[str]]:
    """
    Main preprocessing function that handles the entire presentation preprocessing pipeline.
    """
    print(f"\nStarting presentation preprocessing: {file["name"]}")
    
    print("\nStep 1: Extracting text from presentation")
    content_parts = extract_text_from_pptx(file)
    print(f"Extracted {len(content_parts)} content parts from presentation")
    
    # Process each part according to its type
    processed_parts = []
    for content_type, content in content_parts:
        if content_type == 'table':
            # Keep tables as is
            processed_parts.append(f"\nTable Content:\n{content}")
        else:
            # Clean non-table content
            processed_parts.append(clean_text(content))
    
    # Join all parts together
    final_text = '\n'.join(processed_parts)
    # final_text = final_text.replace("JPEG 2000","")
    print(f"Extracted {len(final_text)} characters of processed text")
    print('Processed text preview:', final_text[:500])
    
    print("\nStep 2: Splitting into chunks")
    chunks = split_text_into_chunks(final_text, chunk_size, overlap)
    
    # Normalize chunks
    normalized_chunks = []
    for chunk in chunks:
        normalized_chunk = normalize_chunk(chunk)
        normalized_chunks.append(normalized_chunk)
        print(f"Normalized chunk: {normalized_chunk[:100]}...")
    
    print(f"Chunk length for presentation {file['name']} = {len(normalized_chunks)}")
    
    # Prepare metadata
    metadata = {
        'user_id': 'example_user',
        'client_id': 'example_client',
        'timestamp_in_seconds': time.time(),
        'num_chunks': len(normalized_chunks),
        'total_chars': len(final_text)
    }
    metadata.update(blob_metadata or {})
    metadata["has_been_preprocessed"] = "True"
    
    print("\nPreprocessing completed successfully!")
    return {
        'result': {
            'chunks': normalized_chunks,
            'metadata': metadata
        }
    }

if __name__ == "__main__":
    # Example usage
    try:
        start_time = time.time()
        pptx_path = r"test_data/6_Image compression.pptx"
        print("\n=== Starting Presentation Processing ===")
        
        # Read the file and create the expected dictionary format
        with open(pptx_path, 'rb') as f:
            file_content = f.read()
        
        file_dict = {"name": pptx_path, "content": file_content}
        result = preprocess_presentation(file_dict)
        
        print("\n=== Processing Results ===")
        print(f"Processed presentation: {result['result']['metadata']['user_id']}")
        print(f"Number of chunks: {result['result']['metadata']['num_chunks']}")
        print(f"Total characters: {result['result']['metadata']['total_chars']}")
        print(f"Processing time: {time.time() - start_time:.2f} seconds")
        
        # Print first chunk as example
        if result['result']['chunks']:
            print(f"\nFirst chunk example:\n{result['result']['chunks'][0]}\n")
            
    except Exception as e:
        print(f"\nError processing presentation: {str(e)}")
